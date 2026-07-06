"""Shared test fixtures for OfficeLens tests."""

from __future__ import annotations

import tempfile
from pathlib import Path

import pytest


@pytest.fixture
def tmp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as d:
        yield Path(d)


@pytest.fixture
def sample_docx(tmp_dir):
    """Create a sample DOCX file for testing."""
    from docx import Document

    doc = Document()
    doc.core_properties.title = "Test Document"
    doc.core_properties.author = "Test Author"
    doc.core_properties.subject = "Testing"

    doc.add_heading("Introduction", level=1)
    doc.add_paragraph("This is a test document for OfficeLens.")
    doc.add_paragraph("It contains multiple paragraphs with different styles.")

    doc.add_heading("Section 1", level=2)
    doc.add_paragraph("First paragraph in section one.")
    doc.add_paragraph("Second paragraph in section one.")

    doc.add_heading("Section 2", level=2)
    doc.add_paragraph("Content in section two with more text.")

    # Add a table
    table = doc.add_table(rows=3, cols=3)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(0, 2).text = "City"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "30"
    table.cell(1, 2).text = "New York"
    table.cell(2, 0).text = "Bob"
    table.cell(2, 1).text = "25"
    table.cell(2, 2).text = "London"

    path = tmp_dir / "test.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def sample_xlsx(tmp_dir):
    """Create a sample XLSX file for testing."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Employees"

    ws.append(["Name", "Age", "City", "Department"])
    ws.append(["Alice", 30, "New York", "Engineering"])
    ws.append(["Bob", 25, "London", "Marketing"])
    ws.append(["Charlie", 35, "Paris", "Engineering"])

    # Add another sheet
    ws2 = wb.create_sheet("Summary")
    ws2.append(["Metric", "Value"])
    ws2.append(["Total Employees", 3])
    ws2.append(["Avg Age", 30])

    path = tmp_dir / "data.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_dir):
    """Create a sample PPTX file for testing."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Slide 1 - Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Presentation Title"
    subtitle = slide.placeholders[1]
    subtitle.text = "Subtitle text here"

    # Slide 2 - Content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Content Slide"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "First bullet point"
    p = tf.add_paragraph()
    p.text = "Second bullet point"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Third bullet point"

    # Slide 3 - With table
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3)
    table_shape = slide.shapes.add_table(3, 3, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "Product"
    table.cell(0, 1).text = "Q1"
    table.cell(0, 2).text = "Q2"
    table.cell(1, 0).text = "Widget A"
    table.cell(1, 1).text = "100"
    table.cell(1, 2).text = "150"
    table.cell(2, 0).text = "Widget B"
    table.cell(2, 1).text = "200"
    table.cell(2, 2).text = "250"

    path = tmp_dir / "slides.pptx"
    prs.save(str(path))
    return path
