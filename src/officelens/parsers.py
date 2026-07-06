"""Office document parsers — DOCX, XLSX, PPTX → common document model."""

from __future__ import annotations

from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any


class DocType(Enum):
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    UNKNOWN = "unknown"

    @classmethod
    def from_path(cls, path: str | Path) -> DocType:
        ext = Path(path).suffix.lower().lstrip(".")
        try:
            return cls(ext)
        except ValueError:
            return cls.UNKNOWN


@dataclass
class TableCell:
    text: str
    row: int
    col: int
    bold: bool = False
    italic: bool = False


@dataclass
class Table:
    rows: int
    cols: int
    cells: list[TableCell] = field(default_factory=list)
    header_row: list[str] = field(default_factory=list)

    def to_dict(self) -> list[list[str]]:
        """Convert to 2D list of strings."""
        grid: list[list[str]] = [["" for _ in range(self.cols)] for _ in range(self.rows)]
        for cell in self.cells:
            if cell.row < self.rows and cell.col < self.cols:
                grid[cell.row][cell.col] = cell.text
        return grid


@dataclass
class TextRun:
    text: str
    bold: bool = False
    italic: bool = False
    underline: bool = False
    strikethrough: bool = False
    font_size: float | None = None
    font_name: str | None = None
    color: str | None = None


@dataclass
class Paragraph:
    text: str
    style: str = ""
    runs: list[TextRun] = field(default_factory=list)
    level: int = 0  # heading level (0 = not a heading)


@dataclass
class Slide:
    number: int
    title: str = ""
    notes: str = ""
    paragraphs: list[Paragraph] = field(default_factory=list)
    tables: list[Table] = field(default_factory=list)
    image_count: int = 0


@dataclass
class Worksheet:
    name: str
    rows: int = 0
    cols: int = 0
    data: list[list[Any]] = field(default_factory=list)
    headers: list[str] = field(default_factory=list)
    formulas: dict[str, str] = field(default_factory=dict)  # cell_ref → formula
    dimensions: str = ""


@dataclass
class Document:
    path: str
    doc_type: DocType
    title: str = ""
    author: str = ""
    subject: str = ""
    creator: str = ""
    created: str = ""
    modified: str = ""
    paragraphs: list[Paragraph] = field(default_factory=list)
    tables: list[Table] = field(default_factory=list)
    worksheets: list[Worksheet] = field(default_factory=list)
    slides: list[Slide] = field(default_factory=list)
    metadata: dict[str, str] = field(default_factory=dict)
    word_count: int = 0
    char_count: int = 0
    page_estimate: int = 0
    image_count: int = 0

    def plain_text(self) -> str:
        """Extract all text content as plain text."""
        parts: list[str] = []
        if self.doc_type == DocType.DOCX:
            for p in self.paragraphs:
                parts.append(p.text)
            for t in self.tables:
                for row in t.to_dict():
                    parts.append("\t".join(row))
        elif self.doc_type == DocType.XLSX:
            for ws in self.worksheets:
                parts.append(f"[Sheet: {ws.name}]")
                for row in ws.data:
                    parts.append("\t".join(str(c) if c is not None else "" for c in row))
        elif self.doc_type == DocType.PPTX:
            for s in self.slides:
                parts.append(f"[Slide {s.number}]")
                if s.title:
                    parts.append(s.title)
                for p in s.paragraphs:
                    parts.append(p.text)
                if s.notes:
                    parts.append(f"[Notes: {s.notes}]")
        return "\n".join(parts)


def _count_words(text: str) -> int:
    return len(text.split())


def _count_chars(text: str) -> int:
    return len(text)


def _estimate_pages(word_count: int) -> int:
    """Rough estimate: ~250 words per page."""
    return max(1, (word_count + 249) // 250)


# ---------------------------------------------------------------------------
# DOCX Parser
# ---------------------------------------------------------------------------


def parse_docx(path: str | Path) -> Document:
    """Parse a .docx file into a Document model."""
    from docx import Document as DocxDocument

    path = str(path)
    doc = DocxDocument(path)

    paragraphs: list[Paragraph] = []
    for para in doc.paragraphs:
        runs = []
        for run in para.runs:
            runs.append(
                TextRun(
                    text=run.text,
                    bold=run.bold or False,
                    italic=run.italic or False,
                    underline=run.underline or False,
                    strikethrough=run.font.strike if run.font.strike else False,
                    font_size=float(run.font.size.pt) if run.font.size else None,
                    font_name=run.font.name,
                    color=run.font.color.rgb.__str__() if run.font.color and run.font.color.rgb else None,
                )
            )
        level = 0
        style_name = para.style.name if para.style else ""
        if style_name.startswith("Heading"):
            try:
                level = int(style_name.replace("Heading", "").strip())
            except ValueError:
                level = 1
        paragraphs.append(
            Paragraph(
                text=para.text,
                style=style_name,
                runs=runs,
                level=level,
            )
        )

    tables: list[Table] = []
    for table in doc.tables:
        tbl_rows = len(table.rows)
        tbl_cols = len(table.columns)
        cells: list[TableCell] = []
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                cells.append(
                    TableCell(
                        text=cell.text,
                        row=i,
                        col=j,
                    )
                )
        header_row = [cells[c].text for c in range(tbl_cols)] if tbl_rows > 0 else []
        tables.append(Table(rows=tbl_rows, cols=tbl_cols, cells=cells, header_row=header_row))

    # Count images
    image_count = 0
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_count += 1
    except Exception:
        pass

    # Metadata
    props = doc.core_properties
    all_text = "\n".join(p.text for p in paragraphs)
    wc = _count_words(all_text)
    cc = _count_chars(all_text)

    metadata: dict[str, str] = {}
    for key in [
        "category",
        "comments",
        "content_status",
        "identifier",
        "keywords",
        "language",
        "last_modified_by",
        "revision",
        "version",
    ]:
        val = getattr(props, key, None)
        if val:
            metadata[key] = str(val)

    return Document(
        path=path,
        doc_type=DocType.DOCX,
        title=props.title or "",
        author=props.author or "",
        subject=props.subject or "",
        creator=props.last_modified_by or props.author or "",
        created=props.created.isoformat() if props.created else "",
        modified=props.modified.isoformat() if props.modified else "",
        paragraphs=paragraphs,
        tables=tables,
        metadata=metadata,
        word_count=wc,
        char_count=cc,
        page_estimate=_estimate_pages(wc),
        image_count=image_count,
    )


# ---------------------------------------------------------------------------
# XLSX Parser
# ---------------------------------------------------------------------------


def parse_xlsx(path: str | Path) -> Document:
    """Parse a .xlsx file into a Document model."""
    from openpyxl import load_workbook

    path = str(path)
    wb = load_workbook(path, data_only=True, read_only=True)

    worksheets: list[Worksheet] = []
    all_text_parts: list[str] = []

    for ws_name in wb.sheetnames:
        ws = wb[ws_name]
        data: list[list[Any]] = []
        formulas: dict[str, str] = {}
        max_row = ws.max_row or 0
        max_col = ws.max_column or 0

        for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
            row_data: list[Any] = []
            for cell in row:
                row_data.append(cell.value)
                if cell.value is not None:
                    all_text_parts.append(str(cell.value))
            data.append(row_data)

        # Get headers from first row
        headers = [str(c) if c is not None else "" for c in data[0]] if data else []

        dimensions = ws.dimensions if hasattr(ws, "dimensions") else ""

        worksheets.append(
            Worksheet(
                name=ws_name,
                rows=max_row,
                cols=max_col,
                data=data,
                headers=headers,
                formulas=formulas,
                dimensions=dimensions,
            )
        )

    # Metadata
    props = wb.properties
    metadata: dict[str, str] = {}
    if props.creator:
        metadata["creator"] = props.creator
    if props.lastModifiedBy:
        metadata["lastModifiedBy"] = props.lastModifiedBy
    if props.title:
        metadata["title"] = props.title
    if props.subject:
        metadata["subject"] = props.subject
    if props.keywords:
        metadata["keywords"] = props.keywords
    if props.description:
        metadata["description"] = props.description

    all_text = "\n".join(all_text_parts)
    wc = _count_words(all_text)
    cc = _count_chars(all_text)

    import contextlib

    with contextlib.suppress(Exception):
        wb.close()

    return Document(
        path=path,
        doc_type=DocType.XLSX,
        title=props.title or "",
        author=props.creator or "",
        subject=props.subject or "",
        creator=props.lastModifiedBy or props.creator or "",
        created=props.created.isoformat() if props.created else "",
        modified=props.modified.isoformat() if props.modified else "",
        worksheets=worksheets,
        metadata=metadata,
        word_count=wc,
        char_count=cc,
        page_estimate=0,  # N/A for spreadsheets
    )


# ---------------------------------------------------------------------------
# PPTX Parser
# ---------------------------------------------------------------------------


def parse_pptx(path: str | Path) -> Document:
    """Parse a .pptx file into a Document model."""
    from pptx import Presentation

    path = str(path)
    prs = Presentation(path)

    slides: list[Slide] = []
    total_images = 0

    for i, slide_layout in enumerate(prs.slides, 1):
        slide = Slide(number=i)

        # Extract text from shapes
        for shape in slide_layout.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        runs = []
                        for run in para.runs:
                            runs.append(
                                TextRun(
                                    text=run.text,
                                    bold=run.font.bold or False,
                                    italic=run.font.italic or False,
                                    underline=run.font.underline or False,
                                )
                            )
                        slide.paragraphs.append(
                            Paragraph(
                                text=text,
                                runs=runs,
                            )
                        )

            if shape.has_table:
                tbl = shape.table
                tbl_rows = len(tbl.rows)
                tbl_cols = len(tbl.columns)
                cells: list[TableCell] = []
                for ri, row in enumerate(tbl.rows):
                    for ci, cell in enumerate(row.cells):
                        cells.append(TableCell(text=cell.text, row=ri, col=ci))
                header_row = [cells[c].text for c in range(tbl_cols)] if tbl_rows > 0 else []
                slide.tables.append(Table(rows=tbl_rows, cols=tbl_cols, cells=cells, header_row=header_row))

            if hasattr(shape, "image"):
                total_images += 1

            try:
                if shape.placeholder_format is not None and shape.placeholder_format.idx == 0:  # title placeholder
                    slide.title = shape.text
            except (ValueError, AttributeError):
                pass  # Not a placeholder shape

        # Notes
        if slide_layout.has_notes_slide:
            notes_frame = slide_layout.notes_slide.notes_text_frame
            slide.notes = notes_frame.text.strip() if notes_frame else ""

        slides.append(slide)

    all_text = "\n".join(p.text for s in slides for p in s.paragraphs)
    wc = _count_words(all_text)
    cc = _count_chars(all_text)

    # Metadata from core properties
    props = prs.core_properties
    metadata: dict[str, str] = {}
    if props.category:
        metadata["category"] = props.category
    if props.comments:
        metadata["comments"] = props.comments
    if props.keywords:
        metadata["keywords"] = props.keywords
    if props.revision:
        metadata["revision"] = str(props.revision)
    if props.version:
        metadata["version"] = str(props.version)

    return Document(
        path=path,
        doc_type=DocType.PPTX,
        title=props.title or "",
        author=props.author or "",
        subject=props.subject or "",
        creator=props.last_modified_by or props.author or "",
        created=props.created.isoformat() if props.created else "",
        modified=props.modified.isoformat() if props.modified else "",
        slides=slides,
        metadata=metadata,
        word_count=wc,
        char_count=cc,
        page_estimate=len(slides),
        image_count=total_images,
    )


# ---------------------------------------------------------------------------
# Unified parser
# ---------------------------------------------------------------------------


def parse(path: str | Path) -> Document:
    """Parse any supported Office document."""
    doc_type = DocType.from_path(path)
    if doc_type == DocType.DOCX:
        return parse_docx(path)
    elif doc_type == DocType.XLSX:
        return parse_xlsx(path)
    elif doc_type == DocType.PPTX:
        return parse_pptx(path)
    else:
        raise ValueError(f"Unsupported file type: {path}")
